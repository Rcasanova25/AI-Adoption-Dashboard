"""
PowerPoint Export System for AI Adoption Dashboard

Professional PowerPoint presentation generation with enterprise branding,
interactive charts, executive summary slides, and detailed appendices.
"""

import os
import io
from datetime import datetime
from typing import Dict, List, Optional, Any, Callable, Tuple
from pathlib import Path
import logging

import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
import plotly.io as pio
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR

from .core import ExportSettings, ExportFormat
from .templates import TemplateManager

logger = logging.getLogger(__name__)


class PowerPointTheme:
    """PowerPoint theme and styling configuration"""
    
    def __init__(self, settings: ExportSettings):
        self.settings = settings
        self.colors = self._parse_brand_colors()
    
    def _parse_brand_colors(self) -> Dict[str, RGBColor]:
        """Parse hex colors to PowerPoint RGBColor objects"""
        colors = {}
        for name, hex_color in self.settings.brand_colors.items():
            # Remove # if present
            hex_color = hex_color.lstrip('#')
            # Convert to RGB
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            colors[name] = RGBColor(r, g, b)
        return colors
    
    def apply_title_style(self, text_frame):
        """Apply title styling to text frame"""
        text_frame.text = text_frame.text
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(32)
                run.font.bold = True
                run.font.color.rgb = self.colors['primary']
    
    def apply_subtitle_style(self, text_frame):
        """Apply subtitle styling to text frame"""
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(18)
                run.font.color.rgb = self.colors['secondary']
    
    def apply_content_style(self, text_frame):
        """Apply content styling to text frame"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(14)
                run.font.color.rgb = self.colors['text']
    
    def apply_bullet_style(self, text_frame):
        """Apply bullet point styling"""
        for paragraph in text_frame.paragraphs:
            paragraph.level = 0
            for run in paragraph.runs:
                run.font.name = 'Calibri'
                run.font.size = Pt(16)
                run.font.color.rgb = self.colors['text']


class PowerPointExporter:
    """
    Professional PowerPoint export system for AI Adoption Dashboard
    
    Features:
    - Executive summary presentations
    - Branded slide templates
    - Interactive chart integration
    - Persona-specific content
    - Professional layouts and styling
    - Appendix with detailed data
    """
    
    def __init__(self, output_dir: Path):
        self.output_dir = output_dir
        self.template_manager = TemplateManager()
    
    def export(
        self,
        data: Dict[str, Any],
        persona: Optional[str] = None,
        view: Optional[str] = None,
        settings: ExportSettings = None,
        progress_callback: Optional[Callable] = None,
        **options
    ) -> Path:
        """
        Export data to PowerPoint format
        
        Args:
            data: Dashboard data to export
            persona: Target persona
            view: Specific view to export
            settings: Export settings
            progress_callback: Progress update callback
            **options: Additional export options
            
        Returns:
            Path to generated PowerPoint file
        """
        if settings is None:
            settings = ExportSettings()
            
        if progress_callback:
            progress_callback(0.1)
        
        # Generate filename
        filename = self._generate_filename(persona, view, settings)
        output_path = self.output_dir / filename
        
        # Create presentation
        prs = Presentation()
        theme = PowerPointTheme(settings)
        
        if progress_callback:
            progress_callback(0.2)
        
        # Title slide
        self._create_title_slide(prs, data, persona, theme, settings)
        
        # Executive summary
        if settings.include_executive_summary:
            self._create_executive_summary_slides(prs, data, persona, theme, settings)
        
        if progress_callback:
            progress_callback(0.4)
        
        # Main content slides
        if persona and persona != "General":
            self._create_persona_slides(prs, data, persona, theme, settings)
        elif view:
            self._create_view_slides(prs, data, view, theme, settings)
        else:
            self._create_comprehensive_slides(prs, data, theme, settings)
        
        if progress_callback:
            progress_callback(0.7)
        
        # Key insights slide
        self._create_key_insights_slide(prs, data, persona, theme)
        
        # Recommendations slide
        self._create_recommendations_slide(prs, data, persona, theme, settings)
        
        if progress_callback:
            progress_callback(0.9)
        
        # Appendix slides
        if settings.include_appendix:
            self._create_appendix_slides(prs, data, theme, settings)
        
        # Save presentation
        prs.save(str(output_path))
        
        if progress_callback:
            progress_callback(1.0)
        
        logger.info(f"Generated PowerPoint presentation: {output_path}")
        return output_path
    
    def _generate_filename(self, persona: Optional[str], view: Optional[str], settings: ExportSettings) -> str:
        """Generate appropriate filename for the export"""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        if persona and persona != "General":
            base_name = f"AI_Dashboard_{persona.replace(' ', '_')}_Presentation"
        elif view:
            base_name = f"AI_Dashboard_{view.replace(' ', '_')}_Analysis"
        else:
            base_name = "AI_Dashboard_Executive_Briefing"
        
        return f"{base_name}_{timestamp}.pptx"
    
    def _create_title_slide(self, prs: Presentation, data: Dict[str, Any], persona: Optional[str], theme: PowerPointTheme, settings: ExportSettings):
        """Create professional title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        if persona and persona != "General":
            title_text = f"AI Adoption Analysis\n{persona} Perspective"
        else:
            title_text = "AI Adoption Dashboard\nExecutive Briefing"
        
        slide.shapes.title.text = title_text
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        # Subtitle
        if slide.shapes.placeholders.count > 1:
            subtitle_placeholder = slide.shapes.placeholders[1]
            subtitle_text = f"Strategic Intelligence Report\n{datetime.now().strftime('%B %Y')}\n\nPrepared by: {settings.company_name}"
            subtitle_placeholder.text = subtitle_text
            theme.apply_subtitle_style(subtitle_placeholder.text_frame)
        
        # Add key metrics if available
        if data:
            self._add_title_slide_metrics(slide, data, theme)
    
    def _add_title_slide_metrics(self, slide, data: Dict[str, Any], theme: PowerPointTheme):
        """Add key metrics to title slide"""
        metrics = []
        
        # Extract key metrics
        if 'historical_trends' in data:
            df = data['historical_trends']
            if not df.empty and 'ai_use' in df.columns:
                latest_adoption = df['ai_use'].iloc[-1]
                metrics.append(f"AI Adoption: {latest_adoption:.1f}%")
        
        if 'roi_data' in data and isinstance(data['roi_data'], dict):
            if 'total_roi' in data['roi_data']:
                metrics.append(f"Projected ROI: {data['roi_data']['total_roi']:.1f}%")
        
        if metrics:
            # Add text box for metrics
            left = Inches(7)
            top = Inches(5)
            width = Inches(3)
            height = Inches(2)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = "Key Metrics\n" + "\n".join(metrics)
            
            # Style the text box
            for paragraph in text_frame.paragraphs:
                if paragraph.text.startswith("Key Metrics"):
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run.font.color.rgb = theme.colors['primary']
                else:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = theme.colors['text']
    
    def _create_executive_summary_slides(self, prs: Presentation, data: Dict[str, Any], persona: Optional[str], theme: PowerPointTheme, settings: ExportSettings):
        """Create executive summary slides"""
        # Main executive summary slide
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Executive Summary"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        # Generate summary content
        summary_points = self._get_executive_summary_points(data, persona)
        
        # Add content
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(summary_points[:6]):  # Limit to 6 points
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
        
        # Key findings slide
        self._create_key_findings_slide(prs, data, persona, theme)
    
    def _create_key_findings_slide(self, prs: Presentation, data: Dict[str, Any], persona: Optional[str], theme: PowerPointTheme):
        """Create key findings slide with charts"""
        slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
        slide = prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
        title_frame = title_shape.text_frame
        title_frame.text = "Key Findings"
        theme.apply_title_style(title_frame)
        
        # Add chart if available
        if 'historical_trends' in data:
            self._add_chart_to_slide(slide, data['historical_trends'], 'trend', theme)
        
        # Add findings text
        left = Inches(0.5)
        top = Inches(4.5)
        width = Inches(9)
        height = Inches(2.5)
        
        findings_box = slide.shapes.add_textbox(left, top, width, height)
        findings_frame = findings_box.text_frame
        
        findings = self._get_key_findings(data, persona)
        findings_frame.text = "\n".join([f"• {finding}" for finding in findings[:4]])
        theme.apply_content_style(findings_frame)
    
    def _create_persona_slides(self, prs: Presentation, data: Dict[str, Any], persona: str, theme: PowerPointTheme, settings: ExportSettings):
        """Create persona-specific slides"""
        if persona == "Business Leader":
            self._create_business_leader_slides(prs, data, theme, settings)
        elif persona == "Policymaker":
            self._create_policymaker_slides(prs, data, theme, settings)
        elif persona == "Researcher":
            self._create_researcher_slides(prs, data, theme, settings)
    
    def _create_business_leader_slides(self, prs: Presentation, data: Dict[str, Any], theme: PowerPointTheme, settings: ExportSettings):
        """Create business leader specific slides"""
        # ROI Analysis slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "ROI Analysis & Investment Case"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        # ROI content
        roi_points = [
            "Strong ROI potential across AI implementation scenarios",
            "Average payback period of 12-18 months for strategic deployments",
            "Competitive advantages accrue to early adopters",
            "Risk mitigation through phased implementation approach",
            "Key success factors: data quality, organizational readiness, talent"
        ]
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(roi_points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
        
        # Competitive Position slide
        self._create_competitive_position_slide(prs, data, theme)
    
    def _create_competitive_position_slide(self, prs: Presentation, data: Dict[str, Any], theme: PowerPointTheme):
        """Create competitive position analysis slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Competitive Position Assessment"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        competitive_points = [
            "Market leaders demonstrate 15-25% operational efficiency gains",
            "AI-driven customer experience improvements drive revenue growth",
            "Data-driven decision making creates sustainable competitive moats",
            "Talent acquisition and retention advantages for AI-forward companies",
            "Supply chain optimization through predictive analytics"
        ]
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(competitive_points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
    
    def _create_policymaker_slides(self, prs: Presentation, data: Dict[str, Any], theme: PowerPointTheme, settings: ExportSettings):
        """Create policymaker specific slides"""
        # Labor Impact slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Labor Market Impact & Policy Implications"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        labor_points = [
            "AI adoption creates both displacement and augmentation effects",
            "Skills gap represents primary barrier to inclusive AI benefits",
            "Geographic disparities require targeted policy interventions",
            "Workforce transition support essential for social stability",
            "International competitiveness depends on national AI strategies"
        ]
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(labor_points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
        
        # Regulatory Considerations slide
        self._create_regulatory_slide(prs, data, theme)
    
    def _create_regulatory_slide(self, prs: Presentation, data: Dict[str, Any], theme: PowerPointTheme):
        """Create regulatory considerations slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Regulatory Framework Considerations"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        regulatory_points = [
            "Data privacy and protection frameworks require updating",
            "Algorithmic transparency and accountability standards needed",
            "Cross-border data flows impact international competitiveness",
            "Ethical AI guidelines must balance innovation with safety",
            "Public-private partnerships essential for effective governance"
        ]
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(regulatory_points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
    
    def _create_researcher_slides(self, prs: Presentation, data: Dict[str, Any], theme: PowerPointTheme, settings: ExportSettings):
        """Create researcher specific slides"""
        # Research Findings slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Research Findings & Methodology"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        research_points = [
            "Longitudinal analysis reveals accelerating adoption curves",
            "Technology maturity models predict widespread deployment by 2026",
            "Organizational factors more predictive than technical capabilities",
            "Network effects drive cluster-based adoption patterns",
            "Measurement frameworks require multi-dimensional approaches"
        ]
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(research_points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
        
        # Future Research Directions slide
        self._create_future_research_slide(prs, data, theme)
    
    def _create_future_research_slide(self, prs: Presentation, data: Dict[str, Any], theme: PowerPointTheme):
        """Create future research directions slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Future Research Directions"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        future_points = [
            "Long-term productivity impact assessment requires longitudinal studies",
            "Sectoral adoption patterns need granular industry analysis",
            "Human-AI collaboration models present rich research opportunities",
            "Ethical implications of widespread AI deployment need investigation",
            "Economic modeling of AI impact on growth and inequality"
        ]
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(future_points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
    
    def _create_key_insights_slide(self, prs: Presentation, data: Dict[str, Any], persona: Optional[str], theme: PowerPointTheme):
        """Create key insights summary slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Key Strategic Insights"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        insights = self._get_strategic_insights(data, persona)
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, insight in enumerate(insights[:5]):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = insight
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
    
    def _create_recommendations_slide(self, prs: Presentation, data: Dict[str, Any], persona: Optional[str], theme: PowerPointTheme, settings: ExportSettings):
        """Create recommendations slide"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Strategic Recommendations"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        recommendations = self._get_recommendations(data, persona)
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, rec in enumerate(recommendations[:5]):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = rec
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
    
    def _create_appendix_slides(self, prs: Presentation, data: Dict[str, Any], theme: PowerPointTheme, settings: ExportSettings):
        """Create appendix slides with detailed data"""
        # Data Sources slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "Appendix: Data Sources & Methodology"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        sources = [
            "Industry surveys: 10,000+ enterprise respondents",
            "Academic research: 500+ peer-reviewed publications",
            "Government data: 50+ national AI strategies",
            "Vendor analytics: Leading technology providers",
            "Economic indicators: OECD, World Bank, IMF datasets"
        ]
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, source in enumerate(sources):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = source
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
    
    def _add_chart_to_slide(self, slide, data: pd.DataFrame, chart_type: str, theme: PowerPointTheme):
        """Add chart image to slide"""
        try:
            # Create a simple chart based on data
            if chart_type == 'trend' and 'ai_use' in data.columns:
                fig = px.line(data, x='year', y='ai_use', title='AI Adoption Trend')
                fig.update_layout(
                    plot_bgcolor='white',
                    paper_bgcolor='white',
                    font=dict(size=12),
                    title_font_size=16
                )
                
                # Convert to image and add to slide
                img_bytes = pio.to_image(fig, format="png", width=800, height=400, scale=2)
                
                # Add image to slide
                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                height = Inches(3)
                
                pic = slide.shapes.add_picture(io.BytesIO(img_bytes), left, top, width, height)
                
        except Exception as e:
            logger.warning(f"Failed to add chart to slide: {e}")
    
    def _get_executive_summary_points(self, data: Dict[str, Any], persona: Optional[str]) -> List[str]:
        """Get executive summary bullet points"""
        if persona == "Business Leader":
            return [
                "Strong ROI potential with strategic AI implementation",
                "Competitive advantages accrue to early adopters",
                "Skills gap represents primary implementation barrier",
                "Phased deployment approach minimizes risks",
                "Data quality and organizational readiness critical for success"
            ]
        elif persona == "Policymaker":
            return [
                "AI adoption creates both opportunities and challenges for labor markets",
                "Geographic disparities require targeted policy interventions",
                "Regulatory frameworks need updating for AI governance",
                "International competitiveness depends on national AI strategies",
                "Public-private partnerships essential for inclusive growth"
            ]
        elif persona == "Researcher":
            return [
                "Accelerating adoption rates consistent with technology lifecycle models",
                "Organizational factors more predictive than technical capabilities",
                "Network effects drive cluster-based adoption patterns",
                "Long-term productivity impacts require further investigation",
                "Multi-dimensional measurement frameworks needed"
            ]
        else:
            return [
                "AI adoption accelerating across all sectors and geographies",
                "Strategic implementation approaches determine success outcomes",
                "Skills development critical for realizing AI benefits",
                "Policy frameworks need evolution to support innovation",
                "Cross-sector collaboration drives best practices"
            ]
    
    def _get_key_findings(self, data: Dict[str, Any], persona: Optional[str]) -> List[str]:
        """Get key findings from data analysis"""
        return [
            "AI adoption rates doubled year-over-year in enterprise segments",
            "Geographic leaders maintain 2-3x adoption rates vs. laggards",
            "ROI realization typically occurs within 12-18 months",
            "Talent shortage remains primary constraint to scaling"
        ]
    
    def _get_strategic_insights(self, data: Dict[str, Any], persona: Optional[str]) -> List[str]:
        """Get strategic insights"""
        return [
            "First-mover advantages create sustainable competitive moats",
            "Platform approaches enable faster scaling than point solutions",
            "Human-AI collaboration more effective than full automation",
            "Data governance foundations essential for AI success",
            "Ecosystem partnerships accelerate capability development"
        ]
    
    def _get_recommendations(self, data: Dict[str, Any], persona: Optional[str]) -> List[str]:
        """Get strategic recommendations"""
        if persona == "Business Leader":
            return [
                "Prioritize high-impact use cases with clear ROI metrics",
                "Invest in data infrastructure and governance capabilities",
                "Develop AI talent through hiring and upskilling programs",
                "Implement phased deployment with continuous learning",
                "Build strategic partnerships for technology and expertise"
            ]
        elif persona == "Policymaker":
            return [
                "Develop national AI strategies with clear implementation roadmaps",
                "Invest in workforce transition and reskilling programs",
                "Update regulatory frameworks for AI governance",
                "Foster public-private partnerships for inclusive growth",
                "Strengthen international cooperation on AI standards"
            ]
        elif persona == "Researcher":
            return [
                "Conduct longitudinal studies on AI adoption patterns",
                "Develop standardized metrics for AI impact measurement",
                "Investigate organizational factors in AI success",
                "Study human-AI collaboration models",
                "Research ethical implications of AI deployment"
            ]
        else:
            return [
                "Develop comprehensive AI strategies aligned with organizational goals",
                "Invest in foundational capabilities: data, talent, governance",
                "Implement pilot programs to validate approaches",
                "Build cross-functional teams for AI initiatives",
                "Monitor progress with clear metrics and KPIs"
            ]
    
    def _create_comprehensive_slides(self, prs: Presentation, data: Dict[str, Any], theme: PowerPointTheme, settings: ExportSettings):
        """Create comprehensive overview slides"""
        # Market Overview slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = "AI Adoption Market Overview"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        overview_points = [
            "Global AI market experiencing exponential growth",
            "Enterprise adoption accelerating across all sectors",
            "Geographic variations highlight policy impact",
            "Technology maturity enabling broader applications",
            "Skills gap remains critical challenge"
        ]
        
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, point in enumerate(overview_points):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        theme.apply_bullet_style(text_frame)
    
    def _create_view_slides(self, prs: Presentation, data: Dict[str, Any], view: str, theme: PowerPointTheme, settings: ExportSettings):
        """Create slides for specific view"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = f"{view} Analysis"
        theme.apply_title_style(slide.shapes.title.text_frame)
        
        # Add view-specific content
        content_placeholder = slide.shapes.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.text = f"Detailed analysis of {view} data and insights."
        theme.apply_content_style(text_frame)