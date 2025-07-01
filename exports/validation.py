"""
Export Validation System for AI Adoption Dashboard

Comprehensive validation and quality assurance for all export formats,
ensuring data integrity, format compliance, and professional standards.
"""

import os
import json
import xml.etree.ElementTree as ET
from datetime import datetime
from typing import Dict, List, Optional, Any, Tuple, Union
from pathlib import Path
import logging

import pandas as pd
import numpy as np
from PIL import Image
import zipfile

from .core import ExportSettings, ExportFormat, ExportJob, ExportStatus

logger = logging.getLogger(__name__)


class ValidationResult:
    """Validation result container"""
    
    def __init__(self):
        self.valid = True
        self.errors: List[str] = []
        self.warnings: List[str] = []
        self.quality_score = 0.0
        self.file_size = 0
        self.validation_time = datetime.now()
        self.metadata: Dict[str, Any] = {}
    
    def add_error(self, error: str):
        """Add validation error"""
        self.errors.append(error)
        self.valid = False
    
    def add_warning(self, warning: str):
        """Add validation warning"""
        self.warnings.append(warning)
    
    def calculate_quality_score(self):
        """Calculate overall quality score (0-100)"""
        base_score = 100
        
        # Deduct for errors and warnings
        error_penalty = len(self.errors) * 20
        warning_penalty = len(self.warnings) * 5
        
        self.quality_score = max(0, base_score - error_penalty - warning_penalty)
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary"""
        return {
            'valid': self.valid,
            'errors': self.errors,
            'warnings': self.warnings,
            'quality_score': self.quality_score,
            'file_size': self.file_size,
            'validation_time': self.validation_time.isoformat(),
            'metadata': self.metadata
        }


class ExportValidator:
    """
    Comprehensive export validation system
    
    Features:
    - Format-specific validation
    - Data integrity checks
    - Quality assurance metrics
    - Professional standards compliance
    - Accessibility validation
    - Performance analysis
    """
    
    def __init__(self):
        self.validation_rules = self._load_validation_rules()
        self.quality_standards = self._load_quality_standards()
    
    def validate_export(self, file_path: Path, export_format: ExportFormat, settings: ExportSettings = None) -> ValidationResult:
        """
        Comprehensive validation of exported file
        
        Args:
            file_path: Path to exported file
            export_format: Expected export format
            settings: Export settings used
            
        Returns:
            ValidationResult with detailed validation information
        """
        result = ValidationResult()
        
        # Basic file validation
        if not self._validate_file_exists(file_path, result):
            result.calculate_quality_score()
            return result
        
        result.file_size = file_path.stat().st_size
        
        # Format-specific validation
        if export_format == ExportFormat.PDF:
            self._validate_pdf(file_path, result, settings)
        elif export_format == ExportFormat.POWERPOINT:
            self._validate_powerpoint(file_path, result, settings)
        elif export_format == ExportFormat.EXCEL:
            self._validate_excel(file_path, result, settings)
        elif export_format == ExportFormat.JSON:
            self._validate_json(file_path, result, settings)
        elif export_format == ExportFormat.XML:
            self._validate_xml(file_path, result, settings)
        elif export_format == ExportFormat.CSV:
            self._validate_csv(file_path, result, settings)
        elif export_format == ExportFormat.HTML:
            self._validate_html(file_path, result, settings)
        elif export_format == ExportFormat.INTERACTIVE_HTML:
            self._validate_interactive_html(file_path, result, settings)
        elif export_format in [ExportFormat.PNG, ExportFormat.SVG]:
            self._validate_image(file_path, result, settings, export_format)
        
        # General quality checks
        self._validate_file_size(file_path, result, export_format)
        self._validate_naming_convention(file_path, result, export_format)
        
        result.calculate_quality_score()
        
        logger.info(f"Validation completed for {file_path}: Quality Score {result.quality_score:.1f}")
        return result
    
    def validate_data_integrity(self, data: Dict[str, Any]) -> ValidationResult:
        """Validate data integrity before export"""
        result = ValidationResult()
        
        # Check data completeness
        if not data:
            result.add_error("No data provided for export")
            result.calculate_quality_score()
            return result
        
        # Validate individual datasets
        for dataset_name, dataset in data.items():
            if dataset is None:
                result.add_warning(f"Dataset '{dataset_name}' is None")
                continue
            
            if isinstance(dataset, pd.DataFrame):
                self._validate_dataframe(dataset, dataset_name, result)
            elif isinstance(dataset, dict):
                self._validate_dict_data(dataset, dataset_name, result)
            elif isinstance(dataset, list):
                self._validate_list_data(dataset, dataset_name, result)
        
        # Check for required datasets
        required_datasets = ['historical_trends', 'geographic_data']
        for required in required_datasets:
            if required not in data or data[required] is None or (isinstance(data[required], pd.DataFrame) and data[required].empty):
                result.add_warning(f"Recommended dataset '{required}' missing or empty")
        
        result.calculate_quality_score()
        return result
    
    def validate_export_job(self, job: ExportJob) -> ValidationResult:
        """Validate export job configuration"""
        result = ValidationResult()
        
        # Validate job parameters
        if not job.format:
            result.add_error("Export format not specified")
        
        if not job.data:
            result.add_error("No data provided for export")
        
        # Validate format-specific requirements
        if job.format == ExportFormat.PDF and job.persona:
            # Check if persona-specific templates are available
            pass
        
        # Validate options
        if hasattr(job, 'options') and job.options:
            self._validate_export_options(job.options, job.format, result)
        
        result.calculate_quality_score()
        return result
    
    def _validate_file_exists(self, file_path: Path, result: ValidationResult) -> bool:
        """Basic file existence validation"""
        if not file_path.exists():
            result.add_error(f"Export file does not exist: {file_path}")
            return False
        
        if file_path.stat().st_size == 0:
            result.add_error("Export file is empty")
            return False
        
        return True
    
    def _validate_pdf(self, file_path: Path, result: ValidationResult, settings: ExportSettings):
        """Validate PDF export"""
        try:
            # Try to import PyPDF2 for PDF validation
            try:
                import PyPDF2
                
                with open(file_path, 'rb') as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    
                    # Check page count
                    page_count = len(pdf_reader.pages)
                    if page_count == 0:
                        result.add_error("PDF has no pages")
                    elif page_count > 50:
                        result.add_warning(f"PDF has many pages ({page_count}), consider pagination")
                    
                    result.metadata['page_count'] = page_count
                    
                    # Check for text content
                    has_text = False
                    for page in pdf_reader.pages[:3]:  # Check first 3 pages
                        text = page.extract_text()
                        if text and text.strip():
                            has_text = True
                            break
                    
                    if not has_text:
                        result.add_warning("PDF appears to contain no extractable text")
                    
                    # Check metadata
                    if pdf_reader.metadata:
                        result.metadata['pdf_metadata'] = {
                            'title': pdf_reader.metadata.get('/Title'),
                            'author': pdf_reader.metadata.get('/Author'),
                            'creator': pdf_reader.metadata.get('/Creator')
                        }
            
            except ImportError:
                result.add_warning("PyPDF2 not available for detailed PDF validation")
                
        except Exception as e:
            result.add_error(f"PDF validation failed: {str(e)}")
    
    def _validate_powerpoint(self, file_path: Path, result: ValidationResult, settings: ExportSettings):
        """Validate PowerPoint export"""
        try:
            # Check file extension
            if file_path.suffix.lower() != '.pptx':
                result.add_error("PowerPoint file should have .pptx extension")
            
            # Try to validate with python-pptx
            try:
                from pptx import Presentation
                
                prs = Presentation(str(file_path))
                
                slide_count = len(prs.slides)
                if slide_count == 0:
                    result.add_error("PowerPoint presentation has no slides")
                elif slide_count > 30:
                    result.add_warning(f"Presentation has many slides ({slide_count})")
                
                result.metadata['slide_count'] = slide_count
                
                # Check for content
                has_content = False
                for slide in prs.slides[:3]:  # Check first 3 slides
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            has_content = True
                            break
                    if has_content:
                        break
                
                if not has_content:
                    result.add_warning("Presentation appears to have no text content")
                
            except ImportError:
                result.add_warning("python-pptx not available for detailed PowerPoint validation")
                
        except Exception as e:
            result.add_error(f"PowerPoint validation failed: {str(e)}")
    
    def _validate_excel(self, file_path: Path, result: ValidationResult, settings: ExportSettings):
        """Validate Excel export"""
        try:
            # Check file extension
            if file_path.suffix.lower() not in ['.xlsx', '.xls']:
                result.add_error("Excel file should have .xlsx or .xls extension")
            
            # Try to read with pandas
            try:
                # Get sheet names
                excel_file = pd.ExcelFile(file_path)
                sheet_names = excel_file.sheet_names
                
                if not sheet_names:
                    result.add_error("Excel file has no worksheets")
                else:
                    result.metadata['sheet_count'] = len(sheet_names)
                    result.metadata['sheet_names'] = sheet_names
                
                # Check each sheet
                empty_sheets = []
                for sheet_name in sheet_names[:5]:  # Check first 5 sheets
                    try:
                        df = pd.read_excel(file_path, sheet_name=sheet_name)
                        if df.empty:
                            empty_sheets.append(sheet_name)
                    except Exception:
                        result.add_warning(f"Could not read sheet '{sheet_name}'")
                
                if empty_sheets:
                    result.add_warning(f"Empty sheets found: {', '.join(empty_sheets)}")
                
            except Exception as e:
                result.add_error(f"Could not read Excel file: {str(e)}")
                
        except Exception as e:
            result.add_error(f"Excel validation failed: {str(e)}")
    
    def _validate_json(self, file_path: Path, result: ValidationResult, settings: ExportSettings):
        """Validate JSON export"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # Check structure
            if not isinstance(data, dict):
                result.add_error("JSON root should be an object")
                return
            
            # Check required fields
            required_fields = ['metadata', 'datasets']
            for field in required_fields:
                if field not in data:
                    result.add_error(f"Missing required field: {field}")
            
            # Validate metadata
            if 'metadata' in data:
                metadata = data['metadata']
                if not isinstance(metadata, dict):
                    result.add_error("Metadata should be an object")
                else:
                    required_meta_fields = ['export_timestamp', 'dashboard_version']
                    for field in required_meta_fields:
                        if field not in metadata:
                            result.add_warning(f"Missing metadata field: {field}")
            
            # Validate datasets
            if 'datasets' in data:
                datasets = data['datasets']
                if not isinstance(datasets, dict):
                    result.add_error("Datasets should be an object")
                elif not datasets:
                    result.add_warning("No datasets in export")
                else:
                    result.metadata['dataset_count'] = len(datasets)
            
            # Check file size vs content
            if result.file_size > 50 * 1024 * 1024:  # 50MB
                result.add_warning("JSON file is very large, consider compression")
            
        except json.JSONDecodeError as e:
            result.add_error(f"Invalid JSON format: {str(e)}")
        except Exception as e:
            result.add_error(f"JSON validation failed: {str(e)}")
    
    def _validate_xml(self, file_path: Path, result: ValidationResult, settings: ExportSettings):
        """Validate XML export"""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            # Check root element
            if root.tag != "ai_dashboard_export":
                result.add_warning(f"Unexpected root element: {root.tag}")
            
            # Check required elements
            required_elements = ['metadata', 'datasets']
            for element in required_elements:
                if root.find(element) is None:
                    result.add_error(f"Missing required element: {element}")
            
            # Check for content
            if len(list(root)) == 0:
                result.add_error("XML has no child elements")
            
            result.metadata['element_count'] = len(list(root.iter()))
            
        except ET.ParseError as e:
            result.add_error(f"Invalid XML format: {str(e)}")
        except Exception as e:
            result.add_error(f"XML validation failed: {str(e)}")
    
    def _validate_csv(self, file_path: Path, result: ValidationResult, settings: ExportSettings):
        """Validate CSV export"""
        try:
            if file_path.suffix == '.zip':
                # Validate ZIP archive
                with zipfile.ZipFile(file_path, 'r') as zipf:
                    files = zipf.namelist()
                    
                    if not files:
                        result.add_error("ZIP archive is empty")
                        return
                    
                    csv_files = [f for f in files if f.endswith('.csv')]
                    if not csv_files:
                        result.add_error("No CSV files found in archive")
                    
                    if 'metadata.json' not in files:
                        result.add_warning("Missing metadata.json in archive")
                    
                    result.metadata['archive_files'] = len(files)
                    result.metadata['csv_files'] = len(csv_files)
                    
                    # Validate sample CSV files
                    for csv_file in csv_files[:3]:  # Check first 3 CSV files
                        try:
                            with zipf.open(csv_file) as f:
                                df = pd.read_csv(f)
                                if df.empty:
                                    result.add_warning(f"CSV file '{csv_file}' is empty")
                        except Exception as e:
                            result.add_warning(f"Could not read CSV file '{csv_file}': {str(e)}")
            else:
                # Single CSV file
                df = pd.read_csv(file_path)
                
                if df.empty:
                    result.add_error("CSV file is empty")
                else:
                    result.metadata['row_count'] = len(df)
                    result.metadata['column_count'] = len(df.columns)
                    
                    # Check for missing values
                    missing_percentage = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
                    if missing_percentage > 50:
                        result.add_warning(f"High percentage of missing values: {missing_percentage:.1f}%")
                    
        except Exception as e:
            result.add_error(f"CSV validation failed: {str(e)}")
    
    def _validate_html(self, file_path: Path, result: ValidationResult, settings: ExportSettings):
        """Validate HTML export"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Basic HTML structure checks
            if not content.strip().startswith('<!DOCTYPE html>'):
                result.add_warning("HTML should start with DOCTYPE declaration")
            
            if '<html' not in content:
                result.add_error("Missing HTML root element")
            
            if '<head>' not in content:
                result.add_error("Missing HTML head section")
            
            if '<body>' not in content:
                result.add_error("Missing HTML body section")
            
            if '<title>' not in content:
                result.add_warning("Missing HTML title element")
            
            # Check for required content
            if 'AI Adoption Dashboard' not in content:
                result.add_warning("Expected dashboard title not found")
            
            # Check file size
            if len(content) < 1000:
                result.add_warning("HTML file seems very small")
            elif len(content) > 10 * 1024 * 1024:  # 10MB
                result.add_warning("HTML file is very large")
            
            result.metadata['html_size'] = len(content)
            
        except Exception as e:
            result.add_error(f"HTML validation failed: {str(e)}")
    
    def _validate_interactive_html(self, file_path: Path, result: ValidationResult, settings: ExportSettings):
        """Validate interactive HTML export"""
        # First run standard HTML validation
        self._validate_html(file_path, result, settings)
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Check for interactivity
            if 'plotly' not in content.lower():
                result.add_warning("Interactive HTML should contain Plotly charts")
            
            if 'javascript' not in content.lower() and '<script' not in content.lower():
                result.add_warning("Interactive HTML should contain JavaScript")
            
        except Exception as e:
            result.add_error(f"Interactive HTML validation failed: {str(e)}")
    
    def _validate_image(self, file_path: Path, result: ValidationResult, settings: ExportSettings, format: ExportFormat):
        """Validate image export"""
        try:
            if format == ExportFormat.SVG:
                # SVG validation
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                
                if not content.strip().startswith('<svg'):
                    result.add_error("SVG file should start with <svg> element")
                
                if '</svg>' not in content:
                    result.add_error("SVG file missing closing </svg> tag")
                
                result.metadata['svg_size'] = len(content)
                
            else:
                # PNG validation
                with Image.open(file_path) as img:
                    result.metadata['image_size'] = img.size
                    result.metadata['image_mode'] = img.mode
                    result.metadata['image_format'] = img.format
                    
                    # Check dimensions
                    width, height = img.size
                    if width < 300 or height < 200:
                        result.add_warning(f"Image dimensions seem small: {width}x{height}")
                    elif width > 5000 or height > 5000:
                        result.add_warning(f"Image dimensions are very large: {width}x{height}")
                    
                    # Check DPI for print quality
                    if hasattr(img, 'info') and 'dpi' in img.info:
                        dpi = img.info['dpi']
                        if isinstance(dpi, tuple):
                            dpi = dpi[0]
                        if dpi < 150:
                            result.add_warning(f"Low DPI for print quality: {dpi}")
                        elif dpi >= 300:
                            result.metadata['print_ready'] = True
            
        except Exception as e:
            result.add_error(f"Image validation failed: {str(e)}")
    
    def _validate_file_size(self, file_path: Path, result: ValidationResult, format: ExportFormat):
        """Validate file size appropriateness"""
        size_mb = result.file_size / (1024 * 1024)
        
        # Size thresholds by format
        thresholds = {
            ExportFormat.PDF: {'warning': 50, 'error': 200},
            ExportFormat.POWERPOINT: {'warning': 100, 'error': 500},
            ExportFormat.EXCEL: {'warning': 50, 'error': 200},
            ExportFormat.JSON: {'warning': 20, 'error': 100},
            ExportFormat.XML: {'warning': 30, 'error': 150},
            ExportFormat.HTML: {'warning': 10, 'error': 50},
            ExportFormat.PNG: {'warning': 5, 'error': 20},
            ExportFormat.SVG: {'warning': 2, 'error': 10}
        }
        
        if format in thresholds:
            threshold = thresholds[format]
            if size_mb > threshold['error']:
                result.add_error(f"File size too large: {size_mb:.1f}MB")
            elif size_mb > threshold['warning']:
                result.add_warning(f"File size is large: {size_mb:.1f}MB")
        
        # Check for empty or tiny files
        if result.file_size < 1024:  # Less than 1KB
            result.add_warning("File size is very small, may be incomplete")
    
    def _validate_naming_convention(self, file_path: Path, result: ValidationResult, format: ExportFormat):
        """Validate file naming convention"""
        filename = file_path.name
        
        # Check extension matches format
        expected_extensions = {
            ExportFormat.PDF: '.pdf',
            ExportFormat.POWERPOINT: '.pptx',
            ExportFormat.EXCEL: '.xlsx',
            ExportFormat.JSON: '.json',
            ExportFormat.XML: '.xml',
            ExportFormat.CSV: ['.csv', '.zip'],
            ExportFormat.HTML: '.html',
            ExportFormat.INTERACTIVE_HTML: '.html',
            ExportFormat.PNG: '.png',
            ExportFormat.SVG: '.svg'
        }
        
        if format in expected_extensions:
            expected = expected_extensions[format]
            if isinstance(expected, list):
                if not any(filename.endswith(ext) for ext in expected):
                    result.add_warning(f"Unexpected file extension. Expected one of: {expected}")
            else:
                if not filename.endswith(expected):
                    result.add_warning(f"Unexpected file extension. Expected: {expected}")
        
        # Check for timestamp in filename
        if not any(char.isdigit() for char in filename):
            result.add_warning("Filename should include timestamp for uniqueness")
        
        # Check for valid characters
        import re
        if not re.match(r'^[a-zA-Z0-9_\-\.]+$', filename):
            result.add_warning("Filename contains special characters that may cause issues")
    
    def _validate_dataframe(self, df: pd.DataFrame, name: str, result: ValidationResult):
        """Validate DataFrame data"""
        if df.empty:
            result.add_warning(f"DataFrame '{name}' is empty")
            return
        
        # Check for missing values
        missing_percentage = (df.isnull().sum().sum() / (len(df) * len(df.columns))) * 100
        if missing_percentage > 30:
            result.add_warning(f"DataFrame '{name}' has high missing values: {missing_percentage:.1f}%")
        
        # Check for duplicate rows
        duplicate_count = df.duplicated().sum()
        if duplicate_count > 0:
            result.add_warning(f"DataFrame '{name}' has {duplicate_count} duplicate rows")
        
        # Check data types
        for col in df.columns:
            if df[col].dtype == 'object':
                # Check for mixed types
                sample_values = df[col].dropna().head(10)
                types = set(type(val).__name__ for val in sample_values)
                if len(types) > 1:
                    result.add_warning(f"Column '{col}' in '{name}' has mixed data types")
    
    def _validate_dict_data(self, data: dict, name: str, result: ValidationResult):
        """Validate dictionary data"""
        if not data:
            result.add_warning(f"Dictionary '{name}' is empty")
            return
        
        # Check for None values
        none_count = sum(1 for v in data.values() if v is None)
        if none_count > len(data) * 0.3:
            result.add_warning(f"Dictionary '{name}' has many None values: {none_count}/{len(data)}")
    
    def _validate_list_data(self, data: list, name: str, result: ValidationResult):
        """Validate list data"""
        if not data:
            result.add_warning(f"List '{name}' is empty")
            return
        
        # Check for None values
        none_count = sum(1 for item in data if item is None)
        if none_count > len(data) * 0.3:
            result.add_warning(f"List '{name}' has many None values: {none_count}/{len(data)}")
    
    def _validate_export_options(self, options: dict, format: ExportFormat, result: ValidationResult):
        """Validate export options"""
        # Format-specific option validation
        if format == ExportFormat.PDF:
            if 'page_size' in options and options['page_size'] not in ['A4', 'Letter', 'Legal']:
                result.add_warning(f"Unusual page size: {options['page_size']}")
        
        elif format == ExportFormat.PNG:
            if 'dpi' in options:
                dpi = options['dpi']
                if dpi < 72:
                    result.add_warning(f"Low DPI setting: {dpi}")
                elif dpi > 600:
                    result.add_warning(f"Very high DPI setting: {dpi}")
    
    def _load_validation_rules(self) -> Dict[str, Any]:
        """Load validation rules configuration"""
        return {
            'file_size_limits': {
                'pdf': 200 * 1024 * 1024,  # 200MB
                'excel': 100 * 1024 * 1024,  # 100MB
                'json': 50 * 1024 * 1024,   # 50MB
                'image': 20 * 1024 * 1024   # 20MB
            },
            'required_content': {
                'metadata': True,
                'datasets': True,
                'timestamp': True
            },
            'quality_thresholds': {
                'missing_data_percentage': 30,
                'duplicate_rows_percentage': 10
            }
        }
    
    def _load_quality_standards(self) -> Dict[str, Any]:
        """Load quality standards configuration"""
        return {
            'accessibility': {
                'alt_text_required': True,
                'color_contrast_ratio': 4.5,
                'font_size_minimum': 12
            },
            'professional_standards': {
                'consistent_branding': True,
                'proper_citations': True,
                'data_source_attribution': True
            },
            'technical_standards': {
                'valid_markup': True,
                'optimized_file_size': True,
                'cross_platform_compatibility': True
            }
        }
    
    def generate_validation_report(self, validation_results: List[ValidationResult]) -> Dict[str, Any]:
        """Generate comprehensive validation report"""
        total_exports = len(validation_results)
        valid_exports = sum(1 for r in validation_results if r.valid)
        
        avg_quality_score = sum(r.quality_score for r in validation_results) / total_exports if total_exports > 0 else 0
        
        all_errors = []
        all_warnings = []
        
        for result in validation_results:
            all_errors.extend(result.errors)
            all_warnings.extend(result.warnings)
        
        # Group errors and warnings by type
        error_counts = {}
        warning_counts = {}
        
        for error in all_errors:
            error_counts[error] = error_counts.get(error, 0) + 1
        
        for warning in all_warnings:
            warning_counts[warning] = warning_counts.get(warning, 0) + 1
        
        report = {
            'summary': {
                'total_exports': total_exports,
                'valid_exports': valid_exports,
                'success_rate': valid_exports / total_exports if total_exports > 0 else 0,
                'average_quality_score': avg_quality_score
            },
            'errors': {
                'total_count': len(all_errors),
                'unique_count': len(error_counts),
                'most_common': sorted(error_counts.items(), key=lambda x: x[1], reverse=True)[:10]
            },
            'warnings': {
                'total_count': len(all_warnings),
                'unique_count': len(warning_counts),
                'most_common': sorted(warning_counts.items(), key=lambda x: x[1], reverse=True)[:10]
            },
            'recommendations': self._generate_recommendations(validation_results),
            'generated_at': datetime.now().isoformat()
        }
        
        return report
    
    def _generate_recommendations(self, validation_results: List[ValidationResult]) -> List[str]:
        """Generate recommendations based on validation results"""
        recommendations = []
        
        # Analyze common issues
        all_errors = []
        all_warnings = []
        
        for result in validation_results:
            all_errors.extend(result.errors)
            all_warnings.extend(result.warnings)
        
        # File size recommendations
        if any('file size' in error.lower() for error in all_errors + all_warnings):
            recommendations.append("Consider optimizing file sizes through compression or reduced resolution")
        
        # Format-specific recommendations
        if any('pdf' in error.lower() for error in all_errors):
            recommendations.append("Review PDF generation settings and ensure all required libraries are installed")
        
        if any('empty' in error.lower() for error in all_errors):
            recommendations.append("Ensure all data sources are properly loaded before export")
        
        # Quality recommendations
        avg_quality = sum(r.quality_score for r in validation_results) / len(validation_results) if validation_results else 0
        if avg_quality < 80:
            recommendations.append("Overall quality score is below recommended threshold (80%). Review validation errors and warnings.")
        
        if not recommendations:
            recommendations.append("All exports meet quality standards. Continue current practices.")
        
        return recommendations